import time
import json
import tempfile
import os
from openai import OpenAI
from django.utils import timezone
from decks.models import Deck, DeckTheme, Flashcard
from .models import AIJob
from langdetect import detect, DetectorFactory
from notifications.signals import ai_deck_ready

DetectorFactory.seed = 0
client = OpenAI()


class AIGenerationService:
    @staticmethod
    def _generate_from_prompt(
        prompt_text: str,
        lang: str = "en",
        mode: str = "subject",
        requested_count: int | None = None,
    ) -> dict:
        """
        Generate deck + flashcards directly from text or extracted content.
        If `requested_count` is provided and > 0, instruct the model to produce exactly that many.
        Otherwise ask for a reasonable number based on the content.
        """

        # sanitize requested_count
        if requested_count is not None:
            try:
                requested_count = int(requested_count)
                if requested_count <= 0:
                    requested_count = None
            except Exception:
                requested_count = None

        count_instruction = (
            f"Generate EXACTLY {requested_count} flashcards (no more, no less)."
            if requested_count
            else "Generate a reasonable number of flashcards that cover the material."
        )

        if mode == "language":
            system_prompt = f"""
You are an AI language tutor helping English speakers learn {lang}.
Given a text sample or topic, create a study deck that helps learners
improve comprehension, vocabulary, and practical usage of {lang}.

{count_instruction}

Generate:
- a descriptive deck title (1 short line)
- a concise 1–2 sentence description explaining the deck's focus
- 3–6 relevant tags

Flashcards:
- Mix directions: some in English asking for {lang} translation, some in {lang} asking for English.
- Include pronunciation, grammar, or usage where relevant.
- Keep questions and answers clear and concise.

Respond strictly in JSON format exactly like this (no extra fields):
{{
  "title": "...",
  "description": "...",
  "tags": ["tag1","tag2"],
  "flashcards": [
    {{"question":"...","answer":"..."}},
    ...
  ]
}}
"""
        else:
            system_prompt = f"""
You are an AI assistant that creates study decks for students and learners.
Language: {lang}

{count_instruction}

Generate:
- a descriptive deck title (1 short line)
- a concise 1–2 sentence description
- 3–6 relevant tags

Respond strictly in JSON format exactly like this (no extra fields):
{{
  "title":"...",
  "description":"...",
  "tags":["tag1","tag2"],
  "flashcards":[
    {{"question":"...","answer":"..."}},
    ...
  ]
}}
"""

        response = client.chat.completions.create(
            model="gpt-5-nano",
            messages=[
                {"role": "system", "content": system_prompt.strip()},
                {"role": "user", "content": prompt_text},
            ],
        )

        try:
            content = response.choices[0].message.content.strip()
            return json.loads(content)
        except Exception as e:
        
            raise ValueError(f"Failed to parse AI response: {e}")

    # -------------------------------
    # File extraction
    # -------------------------------
    @staticmethod
    def _extract_text_from_file(file_obj):
        import mimetypes

        file_type = mimetypes.guess_type(file_obj.name)[0] or ""
        with tempfile.NamedTemporaryFile(delete=False) as tmp:
            tmp.write(file_obj.read())
            tmp_path = tmp.name

        text = ""
        try:
            if "pdf" in file_type:
                from PyPDF2 import PdfReader

                reader = PdfReader(tmp_path)
                for page in reader.pages:
                    text += page.extract_text() or ""
            elif "word" in file_type or file_obj.name.endswith(".docx"):
                from docx import Document

                doc = Document(tmp_path)
                for p in doc.paragraphs:
                    text += p.text + "\n"
            elif "presentation" in file_type or file_obj.name.endswith(".pptx"):
                from pptx import Presentation

                prs = Presentation(tmp_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            elif "text" in file_type or file_obj.name.endswith(".txt"):
                with open(tmp_path, "r", encoding="utf-8") as f:
                    text = f.read()
            else:
                raise ValueError("Unsupported file type.")
        finally:
            try:
                os.remove(tmp_path)
            except Exception:
                pass

        return text[:6000]

    # -------------------------------
    # Image extraction (OCR)
    # -------------------------------
    @staticmethod
    def _extract_text_from_image(image_obj):
        try:
            from PIL import Image
            import pytesseract
        except ImportError:
            raise ValueError("OCR dependencies (pytesseract, Pillow) are not installed.")

        
        pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

        with tempfile.NamedTemporaryFile(delete=False, suffix=".jpg") as tmp:
            tmp.write(image_obj.read())
            tmp_path = tmp.name

        image = Image.open(tmp_path)
        lang_options = "eng+spa+fra+kor+jpn+chi_tra_vert+chi_sim_vert"
        text = pytesseract.image_to_string(image, lang=lang_options)

        try:
            os.remove(tmp_path)
        except Exception:
            pass

        return text[:6000]

    # -------------------------------
    # Main generator
    # -------------------------------
    @classmethod
    def generate_deck(cls, ai_job: AIJob):
        start_time = time.time()
        try:
            ai_job.mark_processing()

            # -------------------------
            # Extract text input
            # -------------------------
            if ai_job.input_type == "prompt":
                text = ai_job.prompt_text or ""
            elif ai_job.input_type == "file":
                text = cls._extract_text_from_file(ai_job.uploaded_file)
            elif ai_job.input_type == "image":
                text = cls._extract_text_from_image(ai_job.uploaded_image)
            else:
                raise ValueError("Unsupported input type.")

            # -------------------------
            # Detect mode / language
            # -------------------------
            detected_lang = "en"
            mode = "subject"
            try:
                detected_lang = detect(text) if text.strip() else "en"
                lower_text = text.lower()
                for lang_name in ["french", "spanish", "korean", "japanese", "chinese"]:
                    if f"study {lang_name}" in lower_text or f"learn {lang_name}" in lower_text:
                        mode = "language"
                        detected_lang = lang_name
                        break
            except Exception:
                detected_lang = "en"

            # -------------------------
            # Requested count logic
            # -------------------------
            requested_count = getattr(ai_job, "requested_count", None)
            # treat zero/None/invalid as "no specific request"
            if requested_count is not None:
                try:
                    requested_count = int(requested_count)
                    if requested_count <= 0:
                        requested_count = None
                except Exception:
                    requested_count = None

            # -------------------------
            # Call AI
            # -------------------------
            data = cls._generate_from_prompt(
                prompt_text=text,
                lang=detected_lang,
                mode=mode,
                requested_count=requested_count,
            )

            # ensure expected structure
            flashcards = data.get("flashcards", [])
            if not isinstance(flashcards, list):
                raise ValueError("AI returned invalid 'flashcards' format (expected list).")

            # -------------------------
            # Enforce exact count if requested (server-side)
            # -------------------------
            if requested_count:
                # trim if too many
                if len(flashcards) > requested_count:
                    flashcards = flashcards[:requested_count]

                
                while len(flashcards) < requested_count:
                    flashcards.append(
                        {
                            "question": "Auto-generated placeholder question",
                            "answer": "Auto-generated placeholder answer",
                        }
                    )

            # optional safety cap: prevent huge decks regardless of AI output
            MAX_ALLOWED = 200
            if len(flashcards) > MAX_ALLOWED:
                flashcards = flashcards[:MAX_ALLOWED]

            # -------------------------
            # Create Deck
            # -------------------------
            deck = Deck.objects.create(
                owner=ai_job.user,
                title=data.get("title", "Untitled Deck"),
                description=data.get("description", "") or "",
                tags=",".join(data.get("tags", [])) if data.get("tags") else "",
                is_public=ai_job.is_public,
            )

            # Assign theme (system default if none exists)
            if not deck.theme:
                default_theme = DeckTheme.objects.filter(owner=ai_job.user, is_default=True).first()
                if not default_theme:
                    default_theme = DeckTheme.objects.filter(is_system_theme=True, is_default=True).first()
                if default_theme:
                    deck.theme = default_theme
                    deck.save(update_fields=["theme"])

            # -------------------------
            # Save flashcards
            # -------------------------
            for card in flashcards:
                Flashcard.objects.create(
                    deck=deck,
                    question=(card.get("question") or "").strip(),
                    answer=(card.get("answer") or "").strip(),
                )

            ai_deck_ready.send(sender=cls, recipient=ai_job.user, deck=deck, job=ai_job)

            ai_job.deck = deck
            ai_job.result_count = len(flashcards)
            ai_job.result_data = data
            ai_job.api_cost = 0.01
            ai_job.generation_time_ms = int((time.time() - start_time) * 1000)
            ai_job.finished_at = timezone.now()
            ai_job.status = "success"
            ai_job.save(update_fields=["deck", "result_count", "result_data", "api_cost", "generation_time_ms", "finished_at", "status"])

            return deck

        except Exception as e:
            # mark error on job and re-raise
            try:
                ai_job.mark_error(str(e))
            except Exception:
                pass
            raise
